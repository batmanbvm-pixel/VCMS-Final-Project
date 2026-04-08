from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

OUT = Path('/Users/patelpreet/Desktop/mid2_mini_project_viva.pptx')

# Color scheme
PRIMARY_COLOR = RGBColor(11, 79, 108)      # Teal
SECONDARY_COLOR = RGBColor(31, 41, 55)    # Dark
ACCENT_COLOR = RGBColor(59, 130, 246)     # Blue
TEXT_COLOR = RGBColor(0, 0, 0)
LIGHT_BG = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)

def add_professional_title_slide(prs, main_title, subtitle1, subtitle2):
    """Professional front page with modern design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Left accent bar
    left_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = ACCENT_COLOR
    left_shape.line.color.rgb = ACCENT_COLOR
    
    # Main Title - Progress Report 2
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = main_title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle 1
    sub1_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(8.5), Inches(0.6))
    sub1_frame = sub1_box.text_frame
    sub1_frame.word_wrap = True
    p = sub1_frame.paragraphs[0]
    p.text = subtitle1
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    
    # Subtitle 2
    sub2_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(8.5), Inches(0.5))
    sub2_frame = sub2_box.text_frame
    sub2_frame.word_wrap = True
    p = sub2_frame.paragraphs[0]
    p.text = subtitle2
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = WHITE
    
    # Bottom info section
    info_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(4.5), Inches(0.4))
    info_frame = info_box.text_frame
    p = info_frame.paragraphs[0]
    p.text = "GUIDED BY"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    name_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(4.5), Inches(0.3))
    name_frame = name_box.text_frame
    p = name_frame.paragraphs[0]
    p.text = "Prachi Shah"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Team Members
    team_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(4.5), Inches(0.3))
    team_frame = team_box.text_frame
    p = team_frame.paragraphs[0]
    p.text = "TEAM MEMBERS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    team_names_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.45), Inches(4.5), Inches(0.85))
    team_names_frame = team_names_box.text_frame
    team_names_frame.word_wrap = True
    p = team_names_frame.paragraphs[0]
    p.text = "Priyanshu Kukadiya • Preet Patel • Avish Patel"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE

def add_content_slide(prs, title, content_points):
    """Consistent content slide design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Top bar with accent
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_COLOR
    top_bar.line.color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Content
    y_position = 1.2
    for point in content_points:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_position), Inches(8.8), Inches(0.55))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        p = content_frame.paragraphs[0]
        p.text = f"• {point}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_COLOR
        p.line_spacing = 1.2
        y_position += 0.55

def add_module_slide(prs, module_no, title, objective, submodules):
    """Module slide with consistent design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Top bar
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_COLOR
    top_bar.line.color.rgb = PRIMARY_COLOR
    
    # Module header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    header_frame = header_box.text_frame
    header_frame.word_wrap = True
    p = header_frame.paragraphs[0]
    p.text = f"Module {module_no}: {title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Objective
    obj_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(0.5))
    obj_frame = obj_box.text_frame
    obj_frame.word_wrap = True
    p = obj_frame.paragraphs[0]
    p.text = f"Objective: {objective}"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # Submodules content
    y_position = 1.65
    for idx, sm in enumerate(submodules, start=1):
        # Submodule title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(9), Inches(0.28))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = f"{module_no}.{idx} {sm['title']}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        y_position += 0.32
        
        # Intro text
        intro_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_position), Inches(8.8), Inches(0.22))
        intro_frame = intro_box.text_frame
        intro_frame.word_wrap = True
        p = intro_frame.paragraphs[0]
        p.text = sm['intro']
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = SECONDARY_COLOR
        y_position += 0.28
        
        # Bullet points (limited for space)
        for point in sm['points'][:2]:
            bullet_box = slide.shapes.add_textbox(Inches(1.0), Inches(y_position), Inches(8.5), Inches(0.24))
            bullet_frame = bullet_box.text_frame
            bullet_frame.word_wrap = True
            p = bullet_frame.paragraphs[0]
            p.text = f"• {point}"
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_COLOR
            y_position += 0.24
        
        if idx < len(submodules):
            y_position += 0.08

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Professional Title - Progress Report 2
    add_professional_title_slide(
        prs,
        "Progress Report 2",
        "VCMS - Virtual Clinic Management System",
        "All Six Modules - Viva Ready"
    )
    
    # Slide 2: User Management
    add_module_slide(prs, 1, "User Management Module", 
        "Manages user identity, authentication, and account lifecycle for Patient, Doctor, and Admin users.",
        [
            {
                'title': 'Role Model and Access Governance',
                'intro': 'User governance is implemented as a strict role-based structure across pages, APIs, and data access.',
                'points': [
                    'Three core roles: Patient, Doctor, and Admin with protected frontend routes',
                    'Backend middleware validates user role before allowing access to restricted operations'
                ]
            },
            {
                'title': 'Registration, Login, and Account Operations',
                'intro': 'Account creation and session flows cover the full user lifecycle from onboarding to secure re-login.',
                'points': [
                    'Registration includes name, email, phone, password, and selected role',
                    'Doctor profile stores specialization, consultation fee, and professional data'
                ]
            },
            {
                'title': 'Security and Persistence',
                'intro': 'Security controls are integrated at credential, session, and API levels.',
                'points': [
                    'Passwords are stored as bcrypt hashes, never plain text',
                    'JWT tokens secure authenticated API usage and protected session flow'
                ]
            }
        ]
    )
    
    # Slide 3: Appointment Management
    add_module_slide(prs, 2, "Appointment Management Module",
        "Provides complete appointment scheduling and tracking from booking request to consultation completion.",
        [
            {
                'title': 'Booking Methods and User Flow',
                'intro': 'Appointments can be created through manual and guided interaction modes for better usability.',
                'points': [
                    'Manual flow: patient selects doctor, date, and time slot directly',
                    'Chatbot flow: asks symptoms, specialization preference, date, and time'
                ]
            },
            {
                'title': 'Status Lifecycle and Business Rules',
                'intro': 'Each booking follows a controlled status model backed by validation rules.',
                'points': [
                    'Status pipeline: Booked, Accepted, In Progress, Completed, Rejected, Cancelled',
                    'Slot conflict checks and schedule validation prevent invalid bookings'
                ]
            },
            {
                'title': 'Realtime Integration and Data Model',
                'intro': 'Appointment data drives downstream modules and real-time user awareness.',
                'points': [
                    'Socket-based updates push immediate appointment status changes to active users',
                    'Appointments collection stores patientId, doctorId, symptoms, notes, status'
                ]
            }
        ]
    )
    
    # Slide 4: Video Consultation
    add_module_slide(prs, 3, "Video Consultation Module",
        "Enables secure in-app teleconsultation so doctor and patient can consult in real time without external meeting tools.",
        [
            {
                'title': 'Consultation Session Flow',
                'intro': 'Video sessions are linked with appointment context to preserve clinical continuity.',
                'points': [
                    'Doctor and patient join consultation room mapped to scheduled appointment',
                    'Session lifecycle states: waiting, active, and ended'
                ]
            },
            {
                'title': 'WebRTC and Signaling',
                'intro': 'Media and signaling layers are implemented using modern web real-time standards.',
                'points': [
                    'WebRTC handles peer audio/video connection with ICE candidate exchange',
                    'Socket.IO signaling exchanges offer, answer, and connection-state events'
                ]
            },
            {
                'title': 'UI Controls and Persistence',
                'intro': 'User interface and backend routes maintain full consultation traceability.',
                'points': [
                    'Video UI supports mic/camera controls, join/leave, and connection indicators',
                    'VideoSessions collection stores appointment linkage, participants, roomId'
                ]
            }
        ]
    )
    
    # Slide 5: Prescription Management
    add_module_slide(prs, 4, "Prescription Management Module",
        "Captures treatment outcomes as digital prescriptions linked to patient consultation records.",
        [
            {
                'title': 'Prescription Lifecycle and Usage',
                'intro': 'Prescription records are generated after consultation and remain available for follow-up care.',
                'points': [
                    'Doctors create and issue prescriptions after consultation completion',
                    'Patients access active and historical prescriptions from dashboard'
                ]
            },
            {
                'title': 'Clinical Data Structure',
                'intro': 'Prescription schema is designed for practical medication management and traceability.',
                'points': [
                    'Data includes medicine name, dosage, frequency, duration, and instructions',
                    'Each record linked to doctorId, patientId, and appointmentId for traceability'
                ]
            },
            {
                'title': 'Frontend, Backend, and Database Coverage',
                'intro': 'Prescription workflow is integrated across application layers.',
                'points': [
                    'Frontend form for prescription creation with field validation',
                    'Backend APIs for create, list, get, update, and download operations'
                ]
            }
        ]
    )
    
    # Slide 6: AI Health Advisory
    add_module_slide(prs, 5, "AI Health Advisory Module",
        "Helps improve decision-making using AI-based assistance for symptom analysis and report simplification.",
        [
            {
                'title': 'Symptom Analysis and Doctor Recommendation',
                'intro': 'Based on patient symptoms, the system recommends a suitable doctor specialization.',
                'points': [
                    'AI analyzes patient-reported symptoms to suggest specialization match',
                    'Doctor recommendation feeds into appointment booking flow'
                ]
            },
            {
                'title': 'Report Analysis and AI Summary Generation',
                'intro': 'Doctor notes are converted into short AI-generated summaries; medical reports are simplified.',
                'points': [
                    'Consultation notes converted to patient-friendly summaries via Gemini AI',
                    'Patients can upload medical reports which are simplified for accessibility'
                ]
            },
            {
                'title': 'AI Integration and Limitations',
                'intro': 'AI outputs are informational only and do not replace medical diagnosis.',
                'points': [
                    'Gemini AI integration for natural language processing',
                    'Disclaimers present on all AI-generated content for legal compliance'
                ]
            }
        ]
    )
    
    # Slide 7: Wellness Module
    add_module_slide(prs, 6, "Wellness Module",
        "Focuses on health awareness and preventive care through family health tracking and inherited risk analysis.",
        [
            {
                'title': 'Family Health Tracking',
                'intro': 'Users can enter health details of family members such as parents or grandparents.',
                'points': [
                    'Family member health details entry form with health conditions',
                    'Support for multiple family members with relationship mapping'
                ]
            },
            {
                'title': 'Inherited Risk Analysis',
                'intro': 'The system analyzes information using simple rules to highlight possible inherited health risks.',
                'points': [
                    'Analysis engine highlights inherited health risks for user awareness',
                    'Data processed live and not stored permanently for privacy'
                ]
            },
            {
                'title': 'Guidance and Privacy',
                'intro': 'Module provides guidance only and does not offer medical conclusions.',
                'points': [
                    'Privacy-first data processing with no permanent storage',
                    'Informational disclaimers present throughout the module'
                ]
            }
        ]
    )
    
    # Slide 8: Technology Stack
    add_content_slide(prs, "Technology Stack",
        [
            "Frontend: React with TypeScript, Tailwind CSS for responsive UI",
            "Backend: Node.js with Express.js for API development and server-side logic",
            "Database: MongoDB Atlas for scalable cloud data storage",
            "Authentication: JWT (JSON Web Tokens) with bcrypt password hashing",
            "Real-time Communication: Socket.IO for updates and WebRTC for video calls",
            "AI Integration: Google Gemini AI for symptom analysis and report simplification",
            "OCR Technology: Medical document scanning and text extraction capabilities"
        ]
    )
    
    # Slide 9: Key Features Summary
    add_content_slide(prs, "Key Features Summary",
        [
            "Role-based access control for Patient, Doctor, and Admin users with separate dashboards",
            "Dual appointment booking: manual selection and AI-assisted chatbot booking",
            "Real-time video consultation with WebRTC technology integrated in web application",
            "Digital prescription management with secure storage and patient download access",
            "AI-powered symptom analysis and medical report simplification",
            "Family health tracking with inherited risk highlighting for preventive care",
            "Real-time notifications and chat integration throughout platform for user engagement"
        ]
    )
    
    # Slide 10: Security & Data Protection
    add_content_slide(prs, "Security & Data Protection",
        [
            "Passwords encrypted using bcrypt hashing algorithm with salt rounds",
            "JWT-based authentication for secure API access and session management",
            "Role-based middleware validation on all protected routes and endpoints",
            "Input validation and sanitization on critical endpoints to prevent attacks",
            "MongoDB Atlas with secure credentials and encrypted connections",
            "Audit-friendly account metadata for compliance tracking and governance",
            "HIPAA-compliant data handling for healthcare information security"
        ]
    )
    
    # Slide 11: Patient Journey Workflow
    add_content_slide(prs, "Patient Journey Workflow",
        [
            "Registration: Patient creates account with basic health information and preferences",
            "Doctor Discovery: AI-powered specialization recommendation based on reported symptoms",
            "Appointment Booking: Manual selection or AI-guided chatbot booking process",
            "Video Consultation: Real-time in-app consultation with doctor in secure virtual room",
            "Prescription Management: Doctor issues digital prescription, patient downloads copy",
            "AI Advisory: Doctor consultation notes converted to simplified summary for patient review",
            "Wellness Tracking: Family health tracking with inherited risk analysis and guidance"
        ]
    )
    
    # Slide 12: Thank You / Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Left accent bar
    left_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = ACCENT_COLOR
    left_shape.line.color.rgb = ACCENT_COLOR
    
    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.word_wrap = True
    p = thanks_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtext
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(2))
    sub_frame = sub_box.text_frame
    sub_frame.word_wrap = True
    p = sub_frame.paragraphs[0]
    p.text = "VCMS Project Complete\nReady for Deployment & Viva"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    
    prs.save(OUT)
    print(f"\n✅ CREATED: {OUT}")
    print(f"✅ 12 Professional Slides - Modern Design")
    print(f"✅ Progress Report 2 - Front Page Professional")
    print(f"✅ Consistent Color Scheme Throughout")
    print(f"✅ No Unwanted Elements - Clean Design")
    print(f"✅ Ready for Viva Presentation\n")

if __name__ == '__main__':
    main()
